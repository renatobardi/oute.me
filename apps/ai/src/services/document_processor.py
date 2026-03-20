import asyncio
import io
import logging

logger = logging.getLogger(__name__)

MAX_EXTRACTED_LENGTH = 10000


async def extract_text(file_bytes: bytes, mime_type: str, filename: str) -> str:
    try:
        if mime_type == "application/pdf":
            return await _extract_pdf(file_bytes)
        elif mime_type in (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        ):
            return await _extract_docx(file_bytes)
        elif mime_type in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ):
            return _extract_xlsx(file_bytes)
        elif mime_type == "text/csv":
            return _extract_csv(file_bytes)
        elif mime_type in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ):
            return _extract_pptx(file_bytes)
        elif mime_type.startswith("image/"):
            return await _extract_image(file_bytes, mime_type)
        else:
            return f"[Tipo não suportado: {mime_type}]"
    except Exception:
        logger.exception("Failed to extract text from %s", filename)
        return f"[Erro ao processar {filename}]"


async def _extract_pdf(file_bytes: bytes) -> str:
    """Usa Document AI Layout Parser se configurado, senão PyMuPDF como fallback."""
    from src.config import settings

    if settings.document_ai_processor_id:
        try:
            result = await _extract_with_document_ai(
                file_bytes, "application/pdf", settings.document_ai_processor_id
            )
            if result:
                return result
        except Exception:
            logger.warning("Document AI falhou para PDF, usando PyMuPDF como fallback")

    return _extract_pdf_pymupdf(file_bytes)


async def _extract_docx(file_bytes: bytes) -> str:
    """Usa Document AI Layout Parser se configurado, senão python-docx como fallback."""
    from src.config import settings

    if settings.document_ai_processor_id:
        try:
            result = await _extract_with_document_ai(
                file_bytes,
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                settings.document_ai_processor_id,
            )
            if result:
                return result
        except Exception:
            logger.warning("Document AI falhou para DOCX, usando python-docx como fallback")

    return _extract_docx_local(file_bytes)


async def _extract_with_document_ai(
    file_bytes: bytes, mime_type: str, processor_id: str
) -> str:
    from google.cloud import documentai

    from src.config import settings

    processor_name = (
        f"projects/{settings.gcp_project}/locations/us/processors/{processor_id}"
    )

    def _call() -> str:
        client = documentai.DocumentProcessorServiceClient()
        raw_document = documentai.RawDocument(content=file_bytes, mime_type=mime_type)
        request = documentai.ProcessRequest(
            name=processor_name, raw_document=raw_document
        )
        result = client.process_document(request=request)
        return result.document.text[:MAX_EXTRACTED_LENGTH]

    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _call)


def _extract_pdf_pymupdf(file_bytes: bytes) -> str:
    import pymupdf

    doc = pymupdf.open(stream=file_bytes, filetype="pdf")  # type: ignore[no-untyped-call]
    text = "\n".join(page.get_text() for page in doc)  # type: ignore[attr-defined]
    doc.close()  # type: ignore[no-untyped-call]
    return text[:MAX_EXTRACTED_LENGTH]


def _extract_docx_local(file_bytes: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return text[:MAX_EXTRACTED_LENGTH]


def _extract_xlsx(file_bytes: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        lines.append(f"[Aba: {sheet}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            lines.append("\t".join(cells))
    wb.close()
    text = "\n".join(lines)
    return text[:MAX_EXTRACTED_LENGTH]


def _extract_csv(file_bytes: bytes) -> str:
    import pandas as pd

    df = pd.read_csv(io.BytesIO(file_bytes))
    return str(df.to_string(max_rows=200))[:MAX_EXTRACTED_LENGTH]


def _extract_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.append(shape.text_frame.text)
    text = "\n".join(lines)
    return text[:MAX_EXTRACTED_LENGTH]


async def _extract_image(file_bytes: bytes, mime_type: str) -> str:
    from vertexai.generative_models import GenerativeModel, Part

    prompt = (
        "Extraia e descreva todo o conteúdo relevante desta imagem"
        " para estimativa de um projeto de software."
        " Inclua textos, diagramas, fluxos, e qualquer informação técnica visível."
    )
    image_part = Part.from_data(mime_type=mime_type, data=file_bytes)
    model = GenerativeModel("gemini-2.5-flash-lite")
    response = await model.generate_content_async([prompt, image_part])  # type: ignore[arg-type]
    return (response.text or "")[:MAX_EXTRACTED_LENGTH]


async def extract_from_url(url: str) -> str:
    import httpx
    from bs4 import BeautifulSoup

    async with httpx.AsyncClient(follow_redirects=True, timeout=15.0) as client:
        resp = await client.get(url)
        resp.raise_for_status()

    soup = BeautifulSoup(resp.text, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    text = soup.get_text(separator="\n", strip=True)
    return text[:MAX_EXTRACTED_LENGTH]
